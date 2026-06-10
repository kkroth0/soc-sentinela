"""
tests/test_cti_slide_pptx.py — Testes do gerador de slide PPTX por feed CTI.
"""

import os

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

from core.models import StandardCTINews
from reports.cti_slide_pptx import build_cti_slide, _fmt_iocs, _fmt_cves, _fmt_list, TOKENS


def _make_template(path: str) -> str:
    """Template-fixture: uma caixa de texto por token, com formatação."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    for i, token in enumerate(TOKENS):
        tb = slide.shapes.add_textbox(Inches(1), Inches(0.5 + i * 0.5), Inches(8), Inches(0.5))
        run = tb.text_frame.paragraphs[0].add_run()
        run.text = f"Campo: {token}"
        run.font.size = Pt(14)
        run.font.bold = True
    prs.save(path)
    return path


def _sample_news() -> StandardCTINews:
    return StandardCTINews(
        title="Zero-day no Exchange Server",
        url="https://example.com/exchange",
        source="BleepingComputer",
        layer=2,
        summary="Primeiro parágrafo do resumo.\n\nSegundo parágrafo com o impacto.",
        date="2026-06-10T13:46:00+00:00",
        iocs={"IPs": ["1.2.3.4"], "Hashes": [{"SHA256": "deadbeef"}]},
        score=72,
        cves=[{"cve_id": "CVE-2026-12345", "cvss_score": 9.8, "risk_tag": "CRITICAL"}],
        sectors=["Governo", "Financeiro"],
        countries=["Brasil"],
    )


def test_build_replaces_all_tokens(tmp_path):
    template = _make_template(str(tmp_path / "tpl.pptx"))
    out = str(tmp_path / "out.pptx")
    result = build_cti_slide(_sample_news(), out, template)

    assert result == out
    assert os.path.exists(out)

    prs = Presentation(out)
    text = "\n".join(
        sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame
    )
    # Nenhum marcador residual
    assert "{{" not in text
    # Conteúdo dos campos presente
    for must in ("Zero-day no Exchange Server", "CVE-2026-12345", "1.2.3.4",
                 "Brasil", "Governo", "72/100", "BleepingComputer", "2026-06-10"):
        assert must in text, f"faltou no slide: {must}"


def test_multiline_value_creates_line_breaks(tmp_path):
    """O RESUMO com \\n\\n deve virar quebras de linha (<a:br>), não texto colado."""
    template = _make_template(str(tmp_path / "tpl.pptx"))
    out = str(tmp_path / "out.pptx")
    build_cti_slide(_sample_news(), out, template)

    prs = Presentation(out)
    breaks = sum(
        len(p._p.findall(qn("a:br")))
        for s in prs.slides for sh in s.shapes if sh.has_text_frame
        for p in sh.text_frame.paragraphs
    )
    assert breaks > 0


def test_formatting_preserved(tmp_path):
    """Os runs reescritos devem manter o <a:rPr> (negrito/tamanho do template)."""
    template = _make_template(str(tmp_path / "tpl.pptx"))
    out = str(tmp_path / "out.pptx")
    build_cti_slide(_sample_news(), out, template)

    prs = Presentation(out)
    titulo_para = next(
        p for s in prs.slides for sh in s.shapes if sh.has_text_frame
        for p in sh.text_frame.paragraphs if "Zero-day" in p.text
    )
    run = titulo_para.runs[0]
    assert run.font.bold is True
    assert run.font.size == Pt(14)


def test_missing_template_raises(tmp_path):
    with pytest.raises(FileNotFoundError):
        build_cti_slide(_sample_news(), str(tmp_path / "o.pptx"), str(tmp_path / "nope.pptx"))


def test_field_formatters():
    assert _fmt_cves([]) == "—"
    assert "CVE-2026-1 (CVSS 5.0 · MEDIUM)" == _fmt_cves(
        [{"cve_id": "CVE-2026-1", "cvss_score": 5.0, "risk_tag": "MEDIUM"}]
    )
    assert "Nenhum" in _fmt_iocs("")
    assert "Nenhum" in _fmt_iocs({"IPs": []})
    assert "1.2.3.4" in _fmt_iocs({"IPs": ["1.2.3.4"]})
    assert _fmt_list([]) == "—"
    assert _fmt_list(["a", "b"]) == "a, b"
