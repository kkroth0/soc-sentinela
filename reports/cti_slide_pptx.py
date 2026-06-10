"""
reports/cti_slide_pptx.py — Gera um slide .pptx a partir do template do analista.

Substitui marcadores ({{TITULO}}, {{RESUMO}}, ...) pelo conteúdo de um
``StandardCTINews``, preservando o design do template (cores, fontes, logo,
layout). A substituição é em nível de parágrafo: concatena o texto dos runs
(tolerando tokens divididos pelo PowerPoint em vários runs), troca os tokens e
reescreve preservando o ``<a:rPr>`` do 1º run; '\\n' no valor vira quebra de
linha (``<a:br>``). Cobre text boxes, placeholders, tabelas e grupos.
"""

import os
from copy import deepcopy
from typing import Any, Iterator

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from core.models import StandardCTINews
from core.logger import get_logger

logger = get_logger("reports.cti_slide_pptx")

_EMPTY = "—"

# Marcadores suportados — referência para o template do analista.
TOKENS = (
    "{{TITULO}}", "{{RESUMO}}", "{{IOCS}}", "{{CVES}}", "{{SETORES}}",
    "{{PAISES}}", "{{FONTE}}", "{{DATA}}", "{{LINK}}", "{{SCORE}}",
)


def _is_empty(v: Any) -> bool:
    return not v or "nenhum" in str(v).lower()


def _fmt_iocs(iocs_raw: Any) -> str:
    """IoCs (dict categorizado ou string) como texto plano multi-linha."""
    if isinstance(iocs_raw, dict):
        lines: list[str] = []
        for key, values in iocs_raw.items():
            if _is_empty(values):
                continue
            items = values if isinstance(values, list) else [values]
            vals: list[str] = []
            for v in items:
                if isinstance(v, dict):  # hashes aninhados {algoritmo: valor}
                    vals += [f"{k2}: {v2}" for k2, v2 in v.items() if not _is_empty(v2)]
                elif not _is_empty(v):
                    vals.append(str(v))
            if vals:
                lines.append(f"{key}: " + ", ".join(vals))
        return "\n".join(lines) if lines else "Nenhum IoC identificado"
    s = str(iocs_raw).strip()
    return s if s and "nenhum" not in s.lower() else "Nenhum IoC identificado"


def _fmt_cves(cves: list[dict[str, Any]]) -> str:
    out: list[str] = []
    for c in (cves or [])[:8]:
        cid = str(c.get("cve_id", "")).strip()
        if not cid:
            continue
        extra: list[str] = []
        if c.get("cvss_score") is not None:
            extra.append(f"CVSS {c['cvss_score']}")
        tag = str(c.get("risk_tag", "")).strip()
        if tag:
            extra.append(tag)
        out.append(cid + (f" ({' · '.join(extra)})" if extra else ""))
    return ", ".join(out) if out else _EMPTY


def _fmt_list(items: list[Any]) -> str:
    vals = [str(x).strip() for x in (items or []) if str(x).strip()]
    return ", ".join(vals) if vals else _EMPTY


def _build_tokens(news: StandardCTINews) -> dict[str, str]:
    """Mapa token → valor textual já formatado para o slide."""
    return {
        "{{TITULO}}": news.title or _EMPTY,
        "{{RESUMO}}": news.summary or _EMPTY,
        "{{IOCS}}": _fmt_iocs(news.iocs),
        "{{CVES}}": _fmt_cves(news.cves),
        "{{SETORES}}": _fmt_list(news.sectors),
        "{{PAISES}}": _fmt_list(news.countries),
        "{{FONTE}}": news.source or _EMPTY,
        "{{DATA}}": (news.date or "")[:10] or _EMPTY,
        "{{LINK}}": news.url or _EMPTY,
        "{{SCORE}}": f"{news.score}/100",
    }


def _iter_paragraphs(shapes: Any) -> Iterator[Any]:
    """Percorre parágrafos de text boxes, placeholders, tabelas e grupos."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_paragraphs(shape.shapes)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield from cell.text_frame.paragraphs
        elif shape.has_text_frame:
            yield from shape.text_frame.paragraphs


def _replace_paragraph(paragraph: Any, tokens: dict[str, str]) -> None:
    """Substitui tokens no parágrafo, preservando a fonte do 1º run."""
    runs = paragraph.runs
    if not runs:
        return
    full = "".join(r.text for r in runs)
    if "{{" not in full:
        return

    new_text = full
    for token, value in tokens.items():
        if token in new_text:
            new_text = new_text.replace(token, value)
    if new_text == full:
        return

    # Formatação a preservar: rPr do 1º run.
    rpr = None
    found = runs[0]._r.find(qn("a:rPr"))
    if found is not None:
        rpr = deepcopy(found)

    p = paragraph._p
    # Remove runs/breaks/fields antigos; mantém pPr e endParaRPr.
    for child in list(p):
        if child.tag in (qn("a:r"), qn("a:br"), qn("a:fld")):
            p.remove(child)

    end = p.find(qn("a:endParaRPr"))

    def _add(elem: Any) -> None:
        if end is not None:
            end.addprevious(elem)
        else:
            p.append(elem)

    for i, seg in enumerate(new_text.split("\n")):
        if i > 0:  # quebra de linha entre segmentos
            br = p.makeelement(qn("a:br"), {})
            if rpr is not None:
                br.append(deepcopy(rpr))
            _add(br)
        r = p.makeelement(qn("a:r"), {})
        if rpr is not None:
            r.append(deepcopy(rpr))
        t = p.makeelement(qn("a:t"), {})
        t.text = seg
        r.append(t)
        _add(r)


def build_cti_slide(news: StandardCTINews, out_path: str, template_path: str) -> str:
    """Gera o .pptx em ``out_path`` a partir de ``template_path`` e o retorna."""
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template PPTX não encontrado: {template_path}")

    tokens = _build_tokens(news)
    prs = Presentation(template_path)
    for slide in prs.slides:
        for paragraph in _iter_paragraphs(slide.shapes):
            _replace_paragraph(paragraph, tokens)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    logger.info("Slide PPTX gerado: %s (fonte: %s)", out_path, news.source)
    return out_path
