"""Gera um template-exemplo de slide CTI com os marcadores em data/templates/.

Serve como default funcional do recurso e como modelo para o analista adaptar
ao design dele (basta manter os marcadores {{...}} nas caixas de texto).
Uso: .venv_temp/bin/python scripts/make_pptx_template.py
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import config

DARK = RGBColor(0x1F, 0x2D, 0x3D)
BLUE = RGBColor(0x2E, 0x86, 0xC1)
GREY = RGBColor(0x55, 0x55, 0x55)


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def _set(tf, text, size, color, bold=False):
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tf


def build_template(out_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Faixa de cabeçalho
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = DARK
    band.line.fill.background()
    htf = band.text_frame
    htf.word_wrap = True
    _set(htf, "{{TITULO}}", 24, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    # Metadados
    _set(_box(slide, 0.5, 1.4, 12.3, 0.5),
         "Fonte: {{FONTE}}   |   Data: {{DATA}}   |   Relevância: {{SCORE}}", 12, GREY)

    # Resumo
    _set(_box(slide, 0.5, 2.1, 12.3, 2.4), "{{RESUMO}}", 14, RGBColor(0x22, 0x22, 0x22))

    # CVEs / Setores / Países
    _set(_box(slide, 0.5, 4.6, 12.3, 0.5), "CVEs: {{CVES}}", 12, BLUE, bold=True)
    _set(_box(slide, 0.5, 5.1, 6.0, 0.5), "Setores-alvo: {{SETORES}}", 12, GREY)
    _set(_box(slide, 6.8, 5.1, 6.0, 0.5), "Países-alvo: {{PAISES}}", 12, GREY)

    # IoCs
    _set(_box(slide, 0.5, 5.7, 12.3, 1.1), "IoCs:\n{{IOCS}}", 11, RGBColor(0x44, 0x44, 0x44))

    # Link
    _set(_box(slide, 0.5, 6.9, 12.3, 0.4), "🔗 {{LINK}}", 10, BLUE)

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build_template(config.PPTX_TEMPLATE_PATH)
    print(f"Template-exemplo gerado: {path}")
